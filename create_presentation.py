"""
Script to create a PowerPoint presentation for EduGen AI
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
PRIMARY_COLOR = RGBColor(59, 130, 246)  # Blue
SECONDARY_COLOR = RGBColor(16, 185, 129)  # Green
TEXT_COLOR = RGBColor(31, 41, 55)  # Dark gray

def add_title_slide(title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

def add_content_slide(title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    tf = content_shape.text_frame
    tf.clear()
    
    for i, item in enumerate(content_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)

def add_two_column_slide(title, left_items, right_items):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Remove default placeholder
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_shape:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Add left text box
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i > 0:
            p = left_tf.add_paragraph()
        else:
            p = left_tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)
    
    # Add right text box
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i > 0:
            p = right_tf.add_paragraph()
        else:
            p = right_tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)

# Slide 1: Title
add_title_slide(
    "📚 EduGen AI",
    "AI-Powered Learning Platform for Engineering Students"
)

# Slide 2: Overview
add_content_slide(
    "📌 Overview",
    [
        "Comprehensive AI-powered learning platform for engineering students",
        "Specialized features for arrear students",
        "Personalized, interactive learning experiences",
        "RAG (Retrieval-Augmented Generation) capabilities",
        "Live Website: edugen-ai-zeta.vercel.app"
    ]
)

# Slide 3: Learning Methodology
add_content_slide(
    "🎯 Proven Learning Methodology",
    [
        "3-Step Learning Flow:",
        "  1. Staff Posts Topic/File",
        "  2. Student Reads via AI Chatbot",
        "  3. Student Takes Quiz → Knowledge Mastery",
        "",
        "Scientific Foundation:",
        "• 75% better retention through active recall",
        "• 6% higher attendance in active learning",
        "• 89% more information retained after one week"
    ]
)

# Slide 4: Tech Stack - Frontend
add_content_slide(
    "🛠️ Tech Stack - Frontend",
    [
        "⚛️ React.js 19.1.0 – Modern UI with hooks",
        "🎨 CSS3 – Responsive design",
        "📱 Progressive Web App – Offline capabilities",
        "🎯 React Router DOM – Client-side routing",
        "📊 Canvas Confetti – Interactive feedback"
    ]
)

# Slide 5: Tech Stack - Backend
add_content_slide(
    "🛠️ Tech Stack - Backend",
    [
        "🚀 Node.js + Express – Primary backend",
        "🐍 Python Flask – RAG API server",
        "🧠 Groq API – High-speed inference (Llama 3.3 70b)",
        "🤖 Google Generative AI – Language models",
        "📄 PyPDF2 – PDF parsing",
        "📰 GNews API – Real-time news"
    ]
)

# Slide 6: Database & AI
add_content_slide(
    "💾 Database & AI",
    [
        "🔥 Firebase – Real-time database & auth",
        "🗄️ Firestore – NoSQL document database",
        "🔐 Firebase Auth – Secure authentication",
        "",
        "AI & ML:",
        "• Custom RAG System with vector database",
        "• Hybrid Search (keyword + semantic)",
        "• Local vector indexing for fast retrieval"
    ]
)

# Slide 7: Key Features - Part 1
add_two_column_slide(
    "🌟 Key Features",
    [
        "🤖 Intelligent AI Tutor",
        "• RAG-Powered Answers",
        "• Content Approval Workflow",
        "• Multi-Mode Chatbot",
        "",
        "👥 Staff Dashboard",
        "• Performance Metrics",
        "• Strength/Weakness Analysis",
        "• File Management"
    ],
    [
        "🛡️ Admin Dashboard",
        "• System Oversight",
        "• Data Visualization",
        "",
        "📚 Learning Tools",
        "• Smart Notes Generation",
        "• Adaptive Quiz Creation",
        "• Progress Tracking"
    ]
)

# Slide 8: Key Features - Part 2
add_content_slide(
    "🌟 Additional Features",
    [
        "🎙️ Voice & Audio: Speech recognition, Text-to-Speech",
        "📱 Responsive Design: Desktop, tablet, mobile",
        "🌙 Dark/Light Mode: Customizable interface",
        "📅 Smart Timetable: Exam and class schedules",
        "⏱️ Study Timer: Pomodoro-style with gamified breaks",
        "🎮 Study Break Games: Tic-Tac-Toe, Memory Match, Tricky Cup",
        "🏆 Gamification: Achievements, leaderboards, streaks"
    ]
)

# Slide 9: Architecture
add_content_slide(
    "🏗️ System Architecture",
    [
        "1. User Entry & Authentication (Firebase Auth)",
        "2. Frontend Layer (React PWA)",
        "   • Student Dashboard",
        "   • Staff Dashboard",
        "   • Admin Dashboard",
        "3. Backend Layer (Microservices)",
        "   • Node.js + Express (Primary)",
        "   • Python Flask (RAG Service)",
        "4. Data & Storage (Firestore + Local Storage)",
        "5. AI Integration (Groq, RAG Pipeline)"
    ]
)

# Slide 10: RAG Pipeline
add_content_slide(
    "🔍 RAG Pipeline Process",
    [
        "1. PDF Upload → Text Extraction",
        "2. Vector Chunking → Store in Vector DB",
        "3. User Query → Vector Similarity Search",
        "4. Context Retrieval → Variable Context Window",
        "5. LLM Query → Context + Prompt → Accurate Answer",
        "",
        "Benefits:",
        "• Context-aware answers from uploaded documents",
        "• Syllabus-aligned content generation",
        "• Fast retrieval with local vector indexing"
    ]
)

# Slide 11: Installation
add_content_slide(
    "🚀 Quick Start",
    [
        "Prerequisites:",
        "• Node.js (v20.x+)",
        "• Python 3.10+",
        "• Firebase account",
        "• API keys (Groq/OpenRouter)",
        "",
        "Setup:",
        "1. Clone repository",
        "2. Install dependencies (npm + pip)",
        "3. Configure .env file",
        "4. Run 3 terminals (Frontend, Node Backend, Python RAG)"
    ]
)

# Slide 12: Expected Outcomes
add_content_slide(
    "📈 Expected Outcomes",
    [
        "✅ Improved Semester Pass Rates",
        "   Structured learning path ensures comprehensive coverage",
        "",
        "✅ Enhanced Conceptual Clarity",
        "   AI-guided explanations target individual weak points",
        "",
        "✅ Better Knowledge Retention",
        "   Multi-step approach creates stronger neural pathways",
        "",
        "✅ Reduced Academic Stress",
        "   Gradual, systematic learning prevents cramming"
    ]
)

# Slide 13: Conclusion
add_title_slide(
    "Thank You!",
    "Built with ❤️ for Engineering Students\n\nLive: edugen-ai-zeta.vercel.app"
)

# Save presentation
output_file = "EduGen_AI_Presentation.pptx"
prs.save(output_file)
print(f"Presentation created successfully: {output_file}")

